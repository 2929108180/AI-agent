"""PPT 文档文本提取 — 基于 python-pptx"""

import io

from pptx import Presentation


def extract_pptx(content: bytes) -> str:
    """从 PPT 中提取每张幻灯片的文本内容。"""
    prs = Presentation(io.BytesIO(content))
    slides_text = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        if texts:
            slides_text.append(f"--- 第 {i} 页 ---\n" + "\n".join(texts))
    return "\n\n".join(slides_text)
