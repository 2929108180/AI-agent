"""PPTX 导出服务 — 元素 JSON → SVG → PPTX"""

import tempfile
from pathlib import Path

from loguru import logger
from pptx import Presentation
from pptx.util import Emu

from app.models.schemas import ExportRequest
from app.utils.svg_builder import elements_to_svg


async def export_pptx(request: ExportRequest) -> str:
    """将所有幻灯片的元素数据转换为 SVG 并打包为 PPTX 文件。"""
    logger.info(
        f"\n{'='*50}\n"
        f"📦 PPTX 导出启动\n"
        f"   幻灯片数: {len(request.slides)}\n"
        f"   文件名: {request.filename}\n"
        f"   各页主题: {[s.theme for s in request.slides]}\n"
        f"   各页元素数: {[len(s.elements) for s in request.slides]}\n"
        f"{'='*50}"
    )
    prs = Presentation()
    prs.slide_width = Emu(12192000)   # 16:9 宽 (25.4cm)
    prs.slide_height = Emu(6858000)   # 16:9 高 (14.29cm)
    blank_layout = prs.slide_layouts[6]  # 空白版式

    temp_dir = Path(tempfile.mkdtemp(prefix="ppt_export_"))

    for i, slide_data in enumerate(request.slides):
        # 元素 JSON → SVG
        elements_dicts = [e.model_dump() for e in slide_data.elements]
        svg_code = elements_to_svg(elements_dicts, slide_data.theme)

        # 写入临时 SVG 文件
        svg_path = temp_dir / f"slide_{i + 1}.svg"
        svg_path.write_text(svg_code, encoding="utf-8")

        # 插入幻灯片
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(svg_path),
            Emu(0), Emu(0),
            prs.slide_width, prs.slide_height,
        )

        logger.info(f"📄 导出幻灯片 {i + 1}/{len(request.slides)} → SVG ({len(svg_code)} 字)")

    # 保存 PPTX
    output_path = str(temp_dir / f"{request.filename}.pptx")
    prs.save(output_path)
    logger.info(f"✅ PPTX 导出完成: {output_path} ({len(request.slides)} 页)")

    return output_path
